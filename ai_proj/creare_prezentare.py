from pptx import Presentation
from pptx.util import Inches
import os

def create_presentation():
    prs = Presentation()

    # Layouts:
    # 0: Title Slide
    # 1: Title and Content
    # 5: Title Only
    # You might want to inspect your presentation template for exact layout indices if these don't work perfectly.

    print("Creating Slide 1: Titlu (Impact Vizual)")
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1] # Placeholder for subtitle
    
    title.text = "Călătorie în Univers: Micii Astronauți și Sistemul Solar"
    subtitle.text = "Sunteți gata de decolare? Punem cu toții căștile de astronauți pe cap! 3... 2... 1... Start!"
    # Add a note for image
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: O rachetă veselă care decolează printre stele."

    print("Creating Slide 2: Ce este Sistemul Solar?")
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1] # Placeholder for content
    title.text = "Ce este Sistemul Solar?"
    body.text = "Sistemul Solar este marea familie a Soarelui. Toate planetele se învârt în jurul lui ca într-un carusel uriaș."
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: O vedere de ansamblu cu Soarele în centru și planetele în jurul lui pe „drumuri” (orbite)."

    print("Creating Slide 3: Soarele – Regele Luminii")
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Soarele – Regele Luminii"
    body.text = "Soarele este o stea uriașă și fierbinte. El ne dă lumină și căldură. Fără el, pe Pământ ar fi foarte frig și întuneric!"
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: Un soare mare, portocaliu, cu față zâmbitoare (dar realistă)."

    print("Creating Slide 4: Mercur și Venus (Cele mai apropiate)")
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Mercur și Venus (Cele mai apropiate)"
    body.text = """Mercur e cel mai mic și fuge cel mai repede în jurul Soarelui.
Venus este cea mai fierbinte planetă. Strălucește ca o bijuterie pe cer!"""
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: Mercur (mic și cenușiu) și Venus (strălucitor și galben)."

    print("Creating Slide 5: Pământ – Casa Noastră (Planeta Albastră)")
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Pământ – Casa Noastră (Planeta Albastră)"
    body.text = """Aceasta este casa noastră! E singura planetă unde există apă, flori, animale și oameni. De sus, arată ca o minge albastră și frumoasă.

Interactivitate: Vedeți unde locuim noi?"""
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: Pământul văzut din spațiu (oceane albastre, continente verzi/maro)."

    print("Creating Slide 6: Marte – Planeta Roșie")
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Marte – Planeta Roșie"
    body.text = "Marte e plin de praf roșu și munți înalți. Acolo trimitem roboței (rovere) să exploreze și să facă poze."
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: O planetă roșie-portocalie, cu praf și pietre."

    print("Creating Slide 7: Jupiter – Gigantul Jucăuș")
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Jupiter – Gigantul Jucăuș"
    body.text = "Jupiter este cel mai mare frate din familie! E atât de mare încât toate celelalte planete ar încăpea în interiorul lui. Are și furtuni uriașe care durează ani de zile!"
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: O planetă uriașă cu dungi colorate și o „pată” mare roșie."

    print("Creating Slide 8: Saturn – Planeta cu Inele")
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Saturn – Planeta cu Inele"
    body.text = "Saturn este campionul eleganței. Are inele magice făcute din bucățele de gheață și praf care sclipesc în lumină."
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: Saturn cu inelele sale spectaculoase din gheață și praf."

    print("Creating Slide 9: Uranus și Neptun – Giganții de Gheață")
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Uranus și Neptun – Giganții de Gheață"
    body.text = "Aici e foarte, foarte frig! Uranus și Neptun sunt planete înghețate, unde bat vânturi foarte puternice."
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: Două planete albastre-verzui, foarte îndepărtate."

    print("Creating Slide 10: Luna – Vecina Pământului")
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Luna – Vecina Pământului"
    body.text = "Luna nu e o planetă, e prietena Pământului care ne veghează noaptea. Ea se învârte mereu în jurul nostru."
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: Luna plină de cratere, lângă Pământ."

    print("Creating Slide 11: Cum este să fii Astronaut?")
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Cum este să fii Astronaut?"
    body.text = "În spațiu nu există greutate, așa că astronauții plutesc ca niște baloane! Ei mănâncă mâncare specială din tuburi și poartă costume groase ca să se protejeze."
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: Un astronaut care plutește în spațiu (costum alb, cască)."

    print("Creating Slide 12: Joc de Verificare (Interactiv)")
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Joc de Verificare"
    body.text = """1. Care planetă are inele frumoase? (Saturn)
2. Unde locuim noi? (Pământul)
3. Cine ne dă căldură în fiecare zi? (Soarele)"""
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: 3 imagini (Soarele, Pământul, Saturn)."

    print("Creating Slide 13: Final – Revenirea pe Pământ")
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Misiune îndeplinită! Bine ați revenit acasă, mici exploratori!"
    slide.notes_slide.notes_text_frame.text = "Imagine sugerată: Racheta care aterizează pe iarbă verde."

    output_path = os.path.join(os.getcwd(), "Calatorie_in_Univers.pptx")
    prs.save(output_path)
    print(f"\nPrezentarea a fost creată cu succes la: {output_path}")

if __name__ == "__main__":
    create_presentation()
